import os
from pptx import Presentation

def extract_slides_from_pptx(file_path):
    presentation = Presentation(file_path)
    slides_text = []

    for slide in presentation.slides:
        slide_text = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text.append(run.text)

        slides_text.append(slide_text)
        
    print(slides_text)
    return slides_text

def find_and_read_ppt_files(root_dir):
    pptx_files = []

    for root, dirs, files in os.walk(root_dir):
        for file in files:
            if file.endswith(".pptx"):
                pptx_path = os.path.join(root, file)
                pptx_files.append(pptx_path)

    return pptx_files

if __name__ == "__main__":
    l = extract_slides_from_pptx("ai-for-private-equity.pptx")

    print(len(l))